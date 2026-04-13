#!/usr/bin/env python3
# coding: utf-8
"""
天気図パワーポイント生成スクリプト

使用法:
  python make_pptx.py YYYYMMDDHH [--output ファイル名.pptx]

スライド構成:
  【GSM グループ1】 上段: 500hPa渦度  /  下段: GSM地上気圧          (2x2)
  【GSM グループ2】 上段: GSM Fax57   /  下段: GSM Fax78             (2x2)
  【GSM グループ3】 上段: 300hPa ジェット  /  下段: 850hPa Qベクター  (2x2)
  【GSM グループ4】 GSM 850hPa相当温位 FT=12,24,36,48h を1枚に4配置  (4in1)
  【ECM グループ1】 上段: ECM 500hPa渦度  /  下段: ECM地上気圧        (2x2)
  【ECM グループ2】 上段: ECM Fax57  /  下段: ECM Fax78               (2x2)
  【ECM グループ3】 ECM 850hPa相当温位 FT=12,24,36,48h を1枚に4配置  (4in1)

モード:
  "2x2"  : 1スライドに2FT×上下2種類=4画像
  "4in1" : 1スライドに指定FT(ft_filter)を2×2配置（同一suffix）
"""

import sys
import re
import argparse
from pathlib import Path
from datetime import datetime, timedelta

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# スライドグループ設定
# ============================================================
SLIDE_GROUPS = [
    # ── GSM ──────────────────────────────────────────────
    {
        "label":         "GSM: 500hPa渦度 / 地上気圧",
        "mode":          "2x2",
        "top_suffix":    "500hPa_Height_VORT",
        "bottom_suffix": "GSM_SurfacePressure",
    },
    {
        "label":         "GSM: 500hPa気温 (Fax57) / 700hPa収束+850hPa気温・風 (Fax78)",
        "mode":          "2x2",
        "top_suffix":    "GSM_Fax57",
        "bottom_suffix": "GSM_Fax78",
    },
    {
        "label":         "GSM: 300hPa ジェット / 850hPa Qベクター",
        "mode":          "2x2",
        "top_suffix":    "300hPa_Jet",
        "bottom_suffix": "850hPa_QVec",
    },
    {
        "label":     "GSM: 850hPa 相当温位",
        "mode":      "4in1",
        "suffix":    "GSM_850hPa_EPT",
        "ft_filter": [12, 24, 36, 48],    # FT=0hは除外
    },
    # ── ECM ──────────────────────────────────────────────
    {
        "label":         "ECM: 500hPa渦度 / 地上気圧",
        "mode":          "2x2",
        "top_suffix":    "ECM_500hPa_Height_VORT",
        "bottom_suffix": "ECM_SurfacePressure",
    },
    {
        "label":         "ECM: 500hPa気温 (Fax57) / 700hPa収束+850hPa気温・風 (Fax78)",
        "mode":          "2x2",
        "top_suffix":    "ECM_Fax57",
        "bottom_suffix": "ECM_Fax78",
    },
    {
        "label":     "ECM: 850hPa 相当温位",
        "mode":      "4in1",
        "suffix":    "ECM_850hPa_EPT",
        "ft_filter": [12, 24, 36, 48],
    },
]

# ── スライドサイズ (4:3 標準) ────────────────────────
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(7.5)

# ── レイアウト定数 ────────────────────────────────────
MARGIN_LR  = Inches(0.12)
MARGIN_TOP = Inches(0.38)
MARGIN_BOT = Inches(0.06)
IMG_GAP_H  = Inches(0.06)
IMG_GAP_V  = Inches(0.06)
LABEL_H    = Inches(0.20)

OUTPUT_DIR = Path(__file__).parent / "output"


# ─────────────────────────────────────────────────────────
# ユーティリティ
# ─────────────────────────────────────────────────────────

def find_ft_list(init_time: str, suffix: str) -> list:
    pattern = re.compile(
        rf"^{re.escape(init_time)}_FT(\d+)h_{re.escape(suffix)}\.png$"
    )
    ft_list = []
    for f in OUTPUT_DIR.glob(f"{init_time}_FT*_{suffix}.png"):
        m = pattern.match(f.name)
        if m:
            ft_list.append(int(m.group(1)))
    return sorted(ft_list)


def get_img_path(init_time: str, ft_h: int, suffix: str) -> Path | None:
    p = OUTPUT_DIR / f"{init_time}_FT{ft_h:03d}h_{suffix}.png"
    return p if p.exists() else None


def valid_time_str(init_time: str, ft_h: int) -> str:
    dt = datetime.strptime(init_time, "%Y%m%d%H") + timedelta(hours=ft_h)
    return dt.strftime("%m/%d %HUTC")


def fit_in_cell(img_path: Path, cell_w: Emu, cell_h: Emu):
    """アスペクト比を保ってセル内に収める。(offset_x, offset_y, draw_w, draw_h) を返す"""
    with PILImage.open(str(img_path)) as im:
        iw, ih = im.size
    cell_ratio = cell_w / cell_h
    img_ratio  = iw / ih
    if img_ratio >= cell_ratio:
        draw_w = cell_w
        draw_h = int(cell_w / img_ratio)
    else:
        draw_h = cell_h
        draw_w = int(cell_h * img_ratio)
    offset_x = (cell_w - draw_w) // 2
    offset_y = (cell_h - draw_h) // 2
    return offset_x, offset_y, draw_w, draw_h


# ─────────────────────────────────────────────────────────
# スライド構築 共通部品
# ─────────────────────────────────────────────────────────

def add_header(slide, init_time: str, ft_list: list, group_label: str):
    w = SLIDE_W - MARGIN_LR * 2
    txb = slide.shapes.add_textbox(
        MARGIN_LR, Inches(0.04), w, MARGIN_TOP - Inches(0.04)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = f"【{group_label}】  "
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(20, 20, 130)

    r2 = p.add_run()
    r2.text = f"IT: {init_time}  "
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor(60, 60, 60)

    r3 = p.add_run()
    r3.text = "  ".join(
        f"FT{ft:03d}h ({valid_time_str(init_time, ft)})" for ft in ft_list
    )
    r3.font.size = Pt(9)
    r3.font.color.rgb = RGBColor(110, 110, 110)


def add_cell(slide, img_p: Path | None,
             cell_left: Emu, cell_top: Emu,
             cell_w: Emu, cell_h: Emu,
             label: str):
    """画像セル1つ (アスペクト比維持・中央配置) + 下部ラベル"""
    img_area_h = cell_h - LABEL_H

    if img_p:
        ox, oy, dw, dh = fit_in_cell(img_p, cell_w, img_area_h)
        slide.shapes.add_picture(
            str(img_p),
            cell_left + ox, cell_top + oy,
            dw, dh
        )
    else:
        box = slide.shapes.add_textbox(cell_left, cell_top, cell_w, img_area_h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "（画像なし）"
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(180, 180, 180)

    lbl_top = cell_top + img_area_h
    lbl_box = slide.shapes.add_textbox(cell_left, lbl_top, cell_w, LABEL_H)
    tf = lbl_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(7.5)
    r.font.color.rgb = RGBColor(90, 90, 90)


# ─────────────────────────────────────────────────────────
# モード別スライド生成
# ─────────────────────────────────────────────────────────

def make_slide_2x2(prs, init_time: str, ft1: int, ft2: int | None, group: dict):
    """2FT × 上下2種類 = 4セル (通常モード)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    ft_disp = [ft1] + ([ft2] if ft2 is not None else [])
    add_header(slide, init_time, ft_disp, group["label"])

    area_w = SLIDE_W - MARGIN_LR * 2
    area_h = SLIDE_H - MARGIN_TOP - MARGIN_BOT
    cell_w = (area_w - IMG_GAP_H) // 2
    cell_h = (area_h - IMG_GAP_V) // 2

    entries = [
        (ft1, group["top_suffix"],    0, 0),
        (ft2, group["top_suffix"],    1, 0),
        (ft1, group["bottom_suffix"], 0, 1),
        (ft2, group["bottom_suffix"], 1, 1),
    ]
    for (ft, suffix, col, row) in entries:
        if ft is None:
            continue
        left  = MARGIN_LR + col * (cell_w + IMG_GAP_H)
        top   = MARGIN_TOP + row * (cell_h + IMG_GAP_V)
        label = f"FT{ft:03d}h ({valid_time_str(init_time, ft)})  {suffix}"
        add_cell(slide, get_img_path(init_time, ft, suffix),
                 left, top, cell_w, cell_h, label)


def make_slide_4in1(prs, init_time: str, group: dict) -> bool:
    """
    ft_filter で指定した4FTを1スライドに2×2配置 (850hPa EPT専用モード)
    左上=FT[0], 右上=FT[1], 左下=FT[2], 右下=FT[3]
    """
    suffix    = group["suffix"]
    ft_filter = group.get("ft_filter", [12, 24, 36, 48])

    # ft_filterのうち少なくとも1枚存在する場合に生成
    available = [ft for ft in ft_filter if get_img_path(init_time, ft, suffix)]
    if not available:
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, init_time, available, group["label"])

    area_w = SLIDE_W - MARGIN_LR * 2
    area_h = SLIDE_H - MARGIN_TOP - MARGIN_BOT
    cell_w = (area_w - IMG_GAP_H) // 2
    cell_h = (area_h - IMG_GAP_V) // 2

    # 最大4枚: (col, row) の順
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for idx, ft in enumerate(ft_filter[:4]):
        col, row = positions[idx]
        left  = MARGIN_LR + col * (cell_w + IMG_GAP_H)
        top   = MARGIN_TOP + row * (cell_h + IMG_GAP_V)
        label = f"FT{ft:03d}h ({valid_time_str(init_time, ft)})  {suffix}"
        add_cell(slide, get_img_path(init_time, ft, suffix),
                 left, top, cell_w, cell_h, label)
    return True


# ─────────────────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="天気図PNG → パワーポイント生成",
        epilog="例: python make_pptx.py 2026041200"
    )
    parser.add_argument("init_time", help="初期時刻 YYYYMMDDHH")
    parser.add_argument(
        "--output", help="出力ファイル名（省略時: tenkizu_YYYYMMDDHH.pptx）"
    )
    args = parser.parse_args()

    init_time = args.init_time
    if len(init_time) != 10 or not init_time.isdigit():
        print("エラー: init_time は YYYYMMDDHH の10桁数字で指定してください")
        sys.exit(1)

    out_path = args.output or f"tenkizu_{init_time}.pptx"

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    total_slides = 0

    for group in SLIDE_GROUPS:
        mode = group.get("mode", "2x2")

        if mode == "4in1":
            suffix = group["suffix"]
            if make_slide_4in1(prs, init_time, group):
                ft_filter = group.get("ft_filter", [12, 24, 36, 48])
                print(f"[グループ] 【{group['label']}】  FT: {ft_filter} → 1スライド")
                total_slides += 1
            else:
                print(f"[スキップ] 【{group['label']}】: 対応画像なし")
            continue

        # 2x2 モード
        ft_top = find_ft_list(init_time, group["top_suffix"])
        ft_bot = find_ft_list(init_time, group["bottom_suffix"])
        ft_all = sorted(set(ft_top) | set(ft_bot))

        if not ft_all:
            print(f"[スキップ] 【{group['label']}】: 対応画像なし")
            continue

        print(f"[グループ] 【{group['label']}】  FT: {ft_all}")

        for i in range(0, len(ft_all), 2):
            ft1 = ft_all[i]
            ft2 = ft_all[i + 1] if i + 1 < len(ft_all) else None
            make_slide_2x2(prs, init_time, ft1, ft2, group)
            total_slides += 1

    if total_slides == 0:
        print(f"エラー: init_time={init_time} に対応する画像が output/ にありません")
        sys.exit(1)

    prs.save(out_path)
    print(f"\n完了: {total_slides}スライド → {out_path}")


if __name__ == "__main__":
    main()
