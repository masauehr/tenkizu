#!/usr/bin/env python3
# coding: utf-8
"""
天気図パワーポイント生成スクリプト (その2 — make_pptx.py の補完版)

make_pptx.py で使用しない残りのGSM画像をパワーポイントにまとめる。

使用法:
  python make_pptx2.py YYYYMMDDHH [--output ファイル名.pptx]

スライド構成:
  【グループ1】 2×2配置
    上段: 300hPa ジェット気流
    下段: 500hPa 気温 (Fax57)

  【グループ2】 2×2配置
    上段: 大気不安定域
    下段: 850hPa 相当温位  ※ 存在するFTのみ

  【鉛直断面図】 1×2配置 (1行2列・大画像)
    左/右: 鉛直断面図  ※ 存在するFTのみ
"""

import sys
import re
import argparse
from pathlib import Path
from datetime import datetime, timedelta

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# スライドグループ設定
# mode "2x2": 上段2枚 + 下段2枚  (通常の天気図ペア)
# mode "1x2": 1行2枚             (鉛直断面図など縦長/特殊図)
# ============================================================
SLIDE_GROUPS = [
    {
        "label":         "GSM: 300hPa ジェット / 500hPa 気温 (Fax57)",
        "mode":          "2x2",
        "top_suffix":    "300hPa_Jet",
        "bottom_suffix": "GSM_Fax57",
    },
    {
        "label":         "GSM: 大気不安定域 / 850hPa 相当温位",
        "mode":          "2x2",
        "top_suffix":    "Instability",
        "bottom_suffix": "GSM_850hPa_EPT",
    },
    {
        "label":         "GSM: 鉛直断面図",
        "mode":          "1x2",
        "suffix":        "CrossSection",
    },
]

# ── スライドサイズ (4:3 標準) ────────────────────────
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(7.5)

# ── レイアウト定数 ────────────────────────────────────
MARGIN_LR  = Inches(0.12)
MARGIN_TOP = Inches(0.38)
MARGIN_BOT = Inches(0.06)
IMG_GAP_H  = Inches(0.06)
IMG_GAP_V  = Inches(0.06)
LABEL_H    = Inches(0.20)

OUTPUT_DIR = Path(__file__).parent / "output"


# ─────────────────────────────────────────────────────────
# ユーティリティ
# ─────────────────────────────────────────────────────────

def find_ft_list(init_time: str, suffix: str) -> list:
    pattern = re.compile(
        rf"^{re.escape(init_time)}_FT(\d+)h_{re.escape(suffix)}\.png$"
    )
    ft_list = []
    for f in OUTPUT_DIR.glob(f"{init_time}_FT*_{suffix}.png"):
        m = pattern.match(f.name)
        if m:
            ft_list.append(int(m.group(1)))
    return sorted(ft_list)


def get_img_path(init_time: str, ft_h: int, suffix: str) -> Path | None:
    p = OUTPUT_DIR / f"{init_time}_FT{ft_h:03d}h_{suffix}.png"
    return p if p.exists() else None


def valid_time_str(init_time: str, ft_h: int) -> str:
    dt = datetime.strptime(init_time, "%Y%m%d%H") + timedelta(hours=ft_h)
    return dt.strftime("%m/%d %HUTC")


def fit_in_cell(img_path: Path, cell_w: Emu, cell_h: Emu):
    """アスペクト比を保ってセル内に収める。(offset_x, offset_y, draw_w, draw_h) を返す"""
    with PILImage.open(str(img_path)) as im:
        iw, ih = im.size
    cell_ratio = cell_w / cell_h
    img_ratio  = iw / ih
    if img_ratio >= cell_ratio:
        draw_w = cell_w
        draw_h = int(cell_w / img_ratio)
    else:
        draw_h = cell_h
        draw_w = int(cell_h * img_ratio)
    offset_x = (cell_w - draw_w) // 2
    offset_y = (cell_h - draw_h) // 2
    return offset_x, offset_y, draw_w, draw_h


# ─────────────────────────────────────────────────────────
# スライド構築
# ─────────────────────────────────────────────────────────

def add_header(slide, init_time: str, ft_list: list, group_label: str):
    w = SLIDE_W - MARGIN_LR * 2
    txb = slide.shapes.add_textbox(
        MARGIN_LR, Inches(0.04), w, MARGIN_TOP - Inches(0.04)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = f"【{group_label}】  "
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(20, 20, 130)

    r2 = p.add_run()
    r2.text = f"IT: {init_time}  "
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor(60, 60, 60)

    r3 = p.add_run()
    r3.text = "  ".join(
        f"FT{ft:03d}h ({valid_time_str(init_time, ft)})" for ft in ft_list
    )
    r3.font.size = Pt(9)
    r3.font.color.rgb = RGBColor(110, 110, 110)


def add_cell(slide, img_p: Path | None,
             cell_left: Emu, cell_top: Emu,
             cell_w: Emu, cell_h: Emu,
             label: str):
    """画像セル1つ (アスペクト比維持・中央配置) + 下部ラベル"""
    img_area_h = cell_h - LABEL_H

    if img_p:
        ox, oy, dw, dh = fit_in_cell(img_p, cell_w, img_area_h)
        slide.shapes.add_picture(
            str(img_p),
            cell_left + ox, cell_top + oy,
            dw, dh
        )
    else:
        box = slide.shapes.add_textbox(cell_left, cell_top, cell_w, img_area_h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "（画像なし）"
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(180, 180, 180)

    lbl_top = cell_top + img_area_h
    lbl_box = slide.shapes.add_textbox(cell_left, lbl_top, cell_w, LABEL_H)
    tf = lbl_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(7.5)
    r.font.color.rgb = RGBColor(90, 90, 90)


# ── 2×2 スライド (上段2枚 + 下段2枚) ─────────────────────

def make_slide_2x2(prs, init_time: str, ft1: int, ft2: int | None, group: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    ft_disp = [ft1] + ([ft2] if ft2 is not None else [])
    add_header(slide, init_time, ft_disp, group["label"])

    area_w = SLIDE_W - MARGIN_LR * 2
    area_h = SLIDE_H - MARGIN_TOP - MARGIN_BOT
    cell_w = (area_w - IMG_GAP_H) // 2
    cell_h = (area_h - IMG_GAP_V) // 2

    entries = [
        (ft1, group["top_suffix"],    0, 0),
        (ft2, group["top_suffix"],    1, 0),
        (ft1, group["bottom_suffix"], 0, 1),
        (ft2, group["bottom_suffix"], 1, 1),
    ]
    for (ft, suffix, col, row) in entries:
        if ft is None:
            continue
        left  = MARGIN_LR + col * (cell_w + IMG_GAP_H)
        top   = MARGIN_TOP + row * (cell_h + IMG_GAP_V)
        label = f"FT{ft:03d}h ({valid_time_str(init_time, ft)})  {suffix}"
        add_cell(slide, get_img_path(init_time, ft, suffix),
                 left, top, cell_w, cell_h, label)


# ── 1×2 スライド (1行2列・大画像) ────────────────────────

def make_slide_1x2(prs, init_time: str, ft1: int, ft2: int | None, group: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    ft_disp = [ft1] + ([ft2] if ft2 is not None else [])
    add_header(slide, init_time, ft_disp, group["label"])

    suffix = group["suffix"]
    area_w = SLIDE_W - MARGIN_LR * 2
    area_h = SLIDE_H - MARGIN_TOP - MARGIN_BOT
    cell_w = (area_w - IMG_GAP_H) // 2
    cell_h = area_h   # 1行なので縦全体を使う

    for col, ft in enumerate([ft1, ft2]):
        if ft is None:
            continue
        left  = MARGIN_LR + col * (cell_w + IMG_GAP_H)
        top   = MARGIN_TOP
        label = f"FT{ft:03d}h ({valid_time_str(init_time, ft)})  {suffix}"
        add_cell(slide, get_img_path(init_time, ft, suffix),
                 left, top, cell_w, cell_h, label)


# ─────────────────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="天気図PNG → パワーポイント生成 (補完版)",
        epilog="例: python make_pptx2.py 2026041200"
    )
    parser.add_argument("init_time", help="初期時刻 YYYYMMDDHH")
    parser.add_argument(
        "--output", help="出力ファイル名（省略時: tenkizu2_YYYYMMDDHH.pptx）"
    )
    args = parser.parse_args()

    init_time = args.init_time
    if len(init_time) != 10 or not init_time.isdigit():
        print("エラー: init_time は YYYYMMDDHH の10桁数字で指定してください")
        sys.exit(1)

    out_path = args.output or f"tenkizu2_{init_time}.pptx"

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    total_slides = 0

    for group in SLIDE_GROUPS:
        mode = group["mode"]

        # FTリストを収集
        if mode == "2x2":
            ft_top = find_ft_list(init_time, group["top_suffix"])
            ft_bot = find_ft_list(init_time, group["bottom_suffix"])
            ft_all = sorted(set(ft_top) | set(ft_bot))
        else:  # 1x2
            ft_all = find_ft_list(init_time, group["suffix"])

        if not ft_all:
            print(f"[スキップ] 【{group['label']}】: 対応画像なし")
            continue

        print(f"[グループ] 【{group['label']}】  FT: {ft_all}")

        for i in range(0, len(ft_all), 2):
            ft1 = ft_all[i]
            ft2 = ft_all[i + 1] if i + 1 < len(ft_all) else None

            if mode == "2x2":
                make_slide_2x2(prs, init_time, ft1, ft2, group)
            else:
                make_slide_1x2(prs, init_time, ft1, ft2, group)

            total_slides += 1

    if total_slides == 0:
        print(f"エラー: init_time={init_time} に対応する画像が output/ にありません")
        sys.exit(1)

    prs.save(out_path)
    print(f"\n完了: {total_slides}スライド → {out_path}")


if __name__ == "__main__":
    main()
